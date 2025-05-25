#!/usr/bin/env python3
"""
Advanced Document Comparison Tool
Converts various document formats (PDF, DOCX, XLSX, PPTX) to images, extracts text with positions, and marks visual differences
"""

import sys
import os
import argparse
import fitz  # PyMuPDF
import difflib
from pathlib import Path
from datetime import datetime
from typing import List, Tuple, Dict, Optional, Set
import re
import html
from PIL import Image, ImageDraw, ImageFont
import base64
import io
import tempfile
import shutil

# Document format conversion imports
try:
    from docx import Document
except ImportError:
    Document = None

try:
    import openpyxl
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.lib.utils import ImageReader
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib import colors
except ImportError:
    openpyxl = None
    canvas = None
    
try:
    from pptx import Presentation
except ImportError:
    Presentation = None


class TextBlock:
    """Represents a text block with position and content."""
    
    def __init__(self, text: str, bbox: Tuple[float, float, float, float], page_num: int):
        self.text = text.strip()
        self.bbox = bbox  # (x0, y0, x1, y1)
        self.page_num = page_num
        self.x0, self.y0, self.x1, self.y1 = bbox
        self.center_x = (self.x0 + self.x1) / 2
        self.center_y = (self.y0 + self.y1) / 2
    
    def __str__(self):
        return f"TextBlock('{self.text[:20]}...', page={self.page_num}, bbox={self.bbox})"


class DocumentComparator:
    """Advanced document comparison with visual annotations. Supports PDF, DOCX, XLSX, PPTX formats."""
    
    def __init__(self):
        self.dpi = 150  # Resolution for document to image conversion
        
        # Use the backend temp directory instead of current working directory
        script_dir = Path(__file__).parent.parent  # Go up to backend directory
        self.temp_dir = script_dir / "temp"
        self.temp_dir.mkdir(exist_ok=True)
        
        # Create timestamped folder for this comparison
        timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
        self.comparison_dir = self.temp_dir / f"compare_{timestamp}"
        self.comparison_dir.mkdir(exist_ok=True)
        
        # Create temp directory for intermediate files
        self.temp_conversion_dir = self.comparison_dir / "temp_conversions"
        self.temp_conversion_dir.mkdir(exist_ok=True)
        
    def get_file_type(self, file_path: str) -> str:
        """Determine the file type based on extension."""
        path = Path(file_path)
        extension = path.suffix.lower()
        
        if extension == '.pdf':
            return 'pdf'
        elif extension in ['.docx', '.doc']:
            return 'docx'
        elif extension in ['.xlsx', '.xls']:
            return 'xlsx'
        elif extension in ['.pptx', '.ppt']:
            return 'pptx'
        else:
            raise ValueError(f"Unsupported file format: {extension}")
    
    def convert_docx_to_pdf(self, docx_path: str) -> str:
        """Convert DOCX file to PDF."""
        if Document is None:
            raise ImportError("python-docx package is required for DOCX support. Install with: pip install python-docx")
        
        print(f"Converting DOCX to PDF: {docx_path}")
        
        # Try using python-docx with reportlab for conversion
        doc = Document(docx_path)
        pdf_path = self.temp_conversion_dir / f"{Path(docx_path).stem}_converted.pdf"
        
        # Create PDF using reportlab
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.pagesizes import letter
        
        doc_pdf = SimpleDocTemplate(str(pdf_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                p = Paragraph(paragraph.text, styles['Normal'])
                story.append(p)
                story.append(Spacer(1, 12))
        
        doc_pdf.build(story)
        return str(pdf_path)
    
    def convert_xlsx_to_pdf(self, xlsx_path: str) -> str:
        """Convert XLSX file to PDF."""
        if openpyxl is None:
            raise ImportError("openpyxl and reportlab packages are required for XLSX support. Install with: pip install openpyxl reportlab")
        
        print(f"Converting XLSX to PDF: {xlsx_path}")
        
        wb = openpyxl.load_workbook(xlsx_path)
        pdf_path = self.temp_conversion_dir / f"{Path(xlsx_path).stem}_converted.pdf"
        
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
        from reportlab.lib.pagesizes import letter, landscape
        from reportlab.lib import colors
        
        doc = SimpleDocTemplate(str(pdf_path), pagesize=landscape(letter))
        story = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            
            # Convert sheet data to list of lists
            data = []
            for row in ws.iter_rows(values_only=True):
                data.append([str(cell) if cell is not None else '' for cell in row])
            
            if data:
                # Create table
                table = Table(data)
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 14),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black)
                ]))
                story.append(table)
        
        doc.build(story)
        return str(pdf_path)
    
    def convert_pptx_to_pdf(self, pptx_path: str) -> str:
        """Convert PPTX file to PDF."""
        if Presentation is None:
            raise ImportError("python-pptx package is required for PPTX support. Install with: pip install python-pptx")
        
        print(f"Converting PPTX to PDF: {pptx_path}")
        
        prs = Presentation(pptx_path)
        pdf_path = self.temp_conversion_dir / f"{Path(pptx_path).stem}_converted.pdf"
        
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.pagesizes import letter
        
        doc = SimpleDocTemplate(str(pdf_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        for i, slide in enumerate(prs.slides):
            story.append(Paragraph(f"Slide {i+1}", styles['Heading1']))
            story.append(Spacer(1, 12))
            
            for shape in slide.shapes:
                try:
                    # Try to get text from different shape types
                    text = ""
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        text = shape.text_frame.text.strip()
                    elif hasattr(shape, 'text'):
                        text = shape.text.strip()
                    
                    if text:
                        p = Paragraph(text, styles['Normal'])
                        story.append(p)
                        story.append(Spacer(1, 6))
                except AttributeError:
                    # Skip shapes that don't have text
                    continue
            
            story.append(Spacer(1, 24))
        
        doc.build(story)
        return str(pdf_path)
    
    def convert_to_pdf(self, file_path: str) -> str:
        """Convert various document formats to PDF."""
        file_type = self.get_file_type(file_path)
        
        if file_type == 'pdf':
            return file_path  # Already PDF
        elif file_type == 'docx':
            return self.convert_docx_to_pdf(file_path)
        elif file_type == 'xlsx':
            return self.convert_xlsx_to_pdf(file_path)
        elif file_type == 'pptx':
            return self.convert_pptx_to_pdf(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
        
    def pdf_to_images(self, pdf_path: str) -> List[Image.Image]:
        """Convert PDF pages to PIL Images."""
        print(f"Converting {pdf_path} to images...")
        doc = fitz.open(pdf_path)
        images = []
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            # Convert to image with high DPI for quality
            mat = fitz.Matrix(self.dpi/72, self.dpi/72)
            pix = page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")
            image = Image.open(io.BytesIO(img_data))
            images.append(image)
            
        doc.close()
        return images
    
    def extract_text_blocks(self, pdf_path: str) -> List[List[TextBlock]]:
        """Extract text blocks with position information from PDF."""
        print(f"Extracting text blocks from {pdf_path}...")
        doc = fitz.open(pdf_path)
        pages_blocks = []
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            blocks = []
            
            # Get text blocks with position
            text_dict = page.get_text("dict")
            for block in text_dict["blocks"]:
                if "lines" in block:  # Text block
                    for line in block["lines"]:
                        line_text = ""
                        line_bbox = line["bbox"]
                        for span in line["spans"]:
                            line_text += span["text"]
                        
                        if line_text.strip():
                            text_block = TextBlock(line_text, line_bbox, page_num)
                            blocks.append(text_block)
            
            pages_blocks.append(blocks)
        
        doc.close()
        return pages_blocks
    
    def normalize_text_for_comparison(self, text: str) -> str:
        """Normalize text for comparison, removing layout-dependent differences."""
        # Remove extra whitespace but preserve word boundaries
        text = re.sub(r'\s+', ' ', text)
        # Remove leading/trailing whitespace
        text = text.strip()
        # Convert to lowercase for case-insensitive comparison
        text = text.lower()
        return text
    
    def find_word_level_differences(self, text1: str, text2: str) -> Dict[str, List]:
        """Find word-level differences between two text strings."""
        words1 = text1.split()
        words2 = text2.split()
        
        matcher = difflib.SequenceMatcher(None, words1, words2)
        word_differences = {
            'deletions': [],    # Word indices deleted from text1
            'insertions': [],   # Word indices added in text2  
            'modifications': [] # Word indices modified between texts
        }
        
        for tag, i1, i2, j1, j2 in matcher.get_opcodes():
            if tag == 'delete':
                word_differences['deletions'].extend(range(i1, i2))
            elif tag == 'insert':
                word_differences['insertions'].extend(range(j1, j2))
            elif tag == 'replace':
                word_differences['modifications'].extend([('old', i) for i in range(i1, i2)])
                word_differences['modifications'].extend([('new', j) for j in range(j1, j2)])
        
        return word_differences

    def find_text_differences(self, blocks1: List[List[TextBlock]], 
                            blocks2: List[List[TextBlock]]) -> Dict:
        """Find text differences between two sets of text blocks with word-level precision."""
        print("Analyzing text differences...")
        
        # Flatten all text blocks and normalize
        text1_lines = []
        text2_lines = []
        block1_map = {}  # normalized_text -> TextBlock
        block2_map = {}  # normalized_text -> TextBlock
        
        for page_blocks in blocks1:
            for block in page_blocks:
                normalized = self.normalize_text_for_comparison(block.text)
                if normalized:
                    text1_lines.append(normalized)
                    block1_map[normalized] = block
        
        for page_blocks in blocks2:
            for block in page_blocks:
                normalized = self.normalize_text_for_comparison(block.text)
                if normalized:
                    text2_lines.append(normalized)
                    block2_map[normalized] = block
        
        # Use difflib to find line-level differences first
        matcher = difflib.SequenceMatcher(None, text1_lines, text2_lines)
        
        differences = {
            'deletions': [],    # Text blocks deleted from pdf1
            'insertions': [],   # Text blocks added in pdf2
            'modifications': [], # Text blocks modified between pdfs
            'word_level': {}    # Map block -> word-level differences
        }
        
        for tag, i1, i2, j1, j2 in matcher.get_opcodes():
            if tag == 'delete':
                # Text deleted from pdf1
                for i in range(i1, i2):
                    if i < len(text1_lines):
                        block = block1_map.get(text1_lines[i])
                        if block:
                            differences['deletions'].append(block)
            
            elif tag == 'insert':
                # Text inserted in pdf2
                for j in range(j1, j2):
                    if j < len(text2_lines):
                        block = block2_map.get(text2_lines[j])
                        if block:
                            differences['insertions'].append(block)
            
            elif tag == 'replace':
                # For replaced lines, perform word-level analysis
                old_lines = text1_lines[i1:i2]
                new_lines = text2_lines[j1:j2]
                
                # If we have corresponding lines, do word-level comparison
                if len(old_lines) == len(new_lines) == 1:
                    old_block = block1_map.get(old_lines[0])
                    new_block = block2_map.get(new_lines[0])
                    
                    if old_block and new_block:
                        # Perform word-level comparison
                        word_diffs = self.find_word_level_differences(old_lines[0], new_lines[0])
                        differences['word_level'][old_block] = ('old', word_diffs, old_lines[0])
                        differences['word_level'][new_block] = ('new', word_diffs, new_lines[0])
                        differences['modifications'].append(('old', old_block))
                        differences['modifications'].append(('new', new_block))
                        continue
                
                # Fall back to line-level for complex changes
                for i in range(i1, i2):
                    if i < len(text1_lines):
                        block = block1_map.get(text1_lines[i])
                        if block:
                            differences['modifications'].append(('old', block))
                
                for j in range(j1, j2):
                    if j < len(text2_lines):
                        block = block2_map.get(text2_lines[j])
                        if block:
                            differences['modifications'].append(('new', block))
        
        return differences
    
    def pad_images_and_blocks(self, images1: List[Image.Image], images2: List[Image.Image],
                            blocks1: List[List[TextBlock]], blocks2: List[List[TextBlock]]) -> Tuple:
        """Pad the shorter PDF with empty pages."""
        max_pages = max(len(images1), len(images2))
        
        # Create empty page template
        if images1:
            width, height = images1[0].size
        elif images2:
            width, height = images2[0].size
        else:
            width, height = 800, 1000  # Default size
        
        empty_image = Image.new('RGB', (width, height), 'white')
        
        # Pad images1 if needed
        while len(images1) < max_pages:
            images1.append(empty_image.copy())
            blocks1.append([])
        
        # Pad images2 if needed
        while len(images2) < max_pages:
            images2.append(empty_image.copy())
            blocks2.append([])
        
        return images1, images2, blocks1, blocks2
    
    def calculate_word_positions(self, text: str, bbox: Tuple[float, float, float, float]) -> List[Tuple[float, float, float, float]]:
        """Calculate approximate positions for each word in a text block."""
        words = text.split()
        if not words:
            return []
        
        x0, y0, x1, y1 = bbox
        total_chars = len(text)
        word_positions = []
        
        current_pos = 0
        for word in words:
            # Calculate word's relative position in the text
            word_start_ratio = current_pos / total_chars if total_chars > 0 else 0
            word_end_ratio = (current_pos + len(word)) / total_chars if total_chars > 0 else 0
            
            # Map to actual coordinates
            word_x0 = x0 + (x1 - x0) * word_start_ratio
            word_x1 = x0 + (x1 - x0) * word_end_ratio
            
            word_positions.append((word_x0, y0, word_x1, y1))
            current_pos += len(word) + 1  # +1 for space
        
        return word_positions

    def annotate_images(self, images1: List[Image.Image], images2: List[Image.Image],
                       differences: Dict) -> Tuple[List[Image.Image], List[Image.Image]]:
        """Annotate images with colored boxes for differences, with word-level precision."""
        print("Annotating images with differences...")
        
        annotated1 = [img.copy() for img in images1]
        annotated2 = [img.copy() for img in images2]
        
        # Scale factor from PDF points to image pixels
        scale_factor = self.dpi / 72
        
        # Track blocks that have word-level differences to avoid double-annotation
        word_level_blocks = set(differences.get('word_level', {}).keys())
        
        # Annotate deletions (light red background with dark red border for visibility)
        for block in differences['deletions']:
            if block.page_num < len(annotated1) and block not in word_level_blocks:
                img = annotated1[block.page_num]
                draw = ImageDraw.Draw(img)
                
                # Scale bbox to image coordinates
                x0 = block.x0 * scale_factor
                y0 = block.y0 * scale_factor
                x1 = block.x1 * scale_factor
                y1 = block.y1 * scale_factor
                
                # Draw light red background to preserve text visibility
                overlay = Image.new('RGBA', img.size, (255, 255, 255, 0))
                overlay_draw = ImageDraw.Draw(overlay)
                overlay_draw.rectangle([x0, y0, x1, y1], fill=(255, 200, 200, 120))  # Light red
                img.paste(overlay, (0, 0), overlay)
                
                # Draw dark red border
                draw.rectangle([x0-2, y0-2, x1+2, y1+2], outline='darkred', width=3)
        
        # Annotate insertions (green on second PDF)
        for block in differences['insertions']:
            if block.page_num < len(annotated2) and block not in word_level_blocks:
                img = annotated2[block.page_num]
                draw = ImageDraw.Draw(img)
                
                # Scale bbox to image coordinates
                x0 = block.x0 * scale_factor
                y0 = block.y0 * scale_factor
                x1 = block.x1 * scale_factor
                y1 = block.y1 * scale_factor
                
                # Draw green box for insertions
                draw.rectangle([x0-2, y0-2, x1+2, y1+2], outline='green', width=3)
                overlay = Image.new('RGBA', img.size, (255, 255, 255, 0))
                overlay_draw = ImageDraw.Draw(overlay)
                overlay_draw.rectangle([x0, y0, x1, y1], fill=(0, 255, 0, 80))
                img.paste(overlay, (0, 0), overlay)
        
        # Annotate word-level modifications
        for block, (change_type, word_diffs, text) in differences.get('word_level', {}).items():
            if change_type == 'old' and block.page_num < len(annotated1):
                img = annotated1[block.page_num]
                self._annotate_word_level_changes(img, block, word_diffs, text, scale_factor, 'old')
                
            elif change_type == 'new' and block.page_num < len(annotated2):
                img = annotated2[block.page_num]
                self._annotate_word_level_changes(img, block, word_diffs, text, scale_factor, 'new')
        
        # Annotate line-level modifications (orange on both PDFs) - only for non-word-level blocks
        for change_type, block in differences['modifications']:
            if block in word_level_blocks:
                continue  # Skip blocks that already have word-level annotations
                
            if change_type == 'old' and block.page_num < len(annotated1):
                img = annotated1[block.page_num]
                draw = ImageDraw.Draw(img)
                
                x0 = block.x0 * scale_factor
                y0 = block.y0 * scale_factor
                x1 = block.x1 * scale_factor
                y1 = block.y1 * scale_factor
                
                draw.rectangle([x0-2, y0-2, x1+2, y1+2], outline='orange', width=3)
                overlay = Image.new('RGBA', img.size, (255, 255, 255, 0))
                overlay_draw = ImageDraw.Draw(overlay)
                overlay_draw.rectangle([x0, y0, x1, y1], fill=(255, 165, 0, 80))
                img.paste(overlay, (0, 0), overlay)
                
            elif change_type == 'new' and block.page_num < len(annotated2):
                img = annotated2[block.page_num]
                draw = ImageDraw.Draw(img)
                
                x0 = block.x0 * scale_factor
                y0 = block.y0 * scale_factor
                x1 = block.x1 * scale_factor
                y1 = block.y1 * scale_factor
                
                draw.rectangle([x0-2, y0-2, x1+2, y1+2], outline='orange', width=3)
                overlay = Image.new('RGBA', img.size, (255, 255, 255, 0))
                overlay_draw = ImageDraw.Draw(overlay)
                overlay_draw.rectangle([x0, y0, x1, y1], fill=(255, 165, 0, 80))
                img.paste(overlay, (0, 0), overlay)
        
        return annotated1, annotated2
    
    def _annotate_word_level_changes(self, img: Image.Image, block: TextBlock, 
                                   word_diffs: Dict, text: str, scale_factor: float, 
                                   change_type: str):
        """Annotate individual words within a text block."""
        words = text.split()
        word_positions = self.calculate_word_positions(text, block.bbox)
        
        draw = ImageDraw.Draw(img)
        overlay = Image.new('RGBA', img.size, (255, 255, 255, 0))
        overlay_draw = ImageDraw.Draw(overlay)
        
        for word_idx, (word_x0, word_y0, word_x1, word_y1) in enumerate(word_positions):
            # Scale to image coordinates
            x0 = word_x0 * scale_factor
            y0 = word_y0 * scale_factor
            x1 = word_x1 * scale_factor
            y1 = word_y1 * scale_factor
            
            # Check if this word has changes
            is_deleted = word_idx in word_diffs.get('deletions', [])
            is_inserted = word_idx in word_diffs.get('insertions', [])
            is_modified = any(item[1] == word_idx for item in word_diffs.get('modifications', []) 
                            if item[0] == change_type)
            
            if is_deleted:
                # Light red background for deletions to preserve visibility
                overlay_draw.rectangle([x0, y0, x1, y1], fill=(255, 200, 200, 150))
                draw.rectangle([x0-1, y0-1, x1+1, y1+1], outline='darkred', width=2)
                
            elif is_inserted:
                # Green for insertions
                overlay_draw.rectangle([x0, y0, x1, y1], fill=(200, 255, 200, 120))
                draw.rectangle([x0-1, y0-1, x1+1, y1+1], outline='green', width=2)
                
            elif is_modified:
                # Orange for modifications
                overlay_draw.rectangle([x0, y0, x1, y1], fill=(255, 220, 180, 120))
                draw.rectangle([x0-1, y0-1, x1+1, y1+1], outline='orange', width=2)
        
        img.paste(overlay, (0, 0), overlay)
    
    def save_images_to_base64(self, images: List[Image.Image], prefix: str) -> List[str]:
        """Save images and return base64 encoded strings for HTML embedding."""
        base64_images = []
        
        for i, img in enumerate(images):
            # Save to comparison directory
            img_path = self.comparison_dir / f"{prefix}_page_{i+1}.png"
            img.save(img_path)
            
            # Convert to base64 for HTML embedding
            buffer = io.BytesIO()
            img.save(buffer, format='PNG')
            img_data = base64.b64encode(buffer.getvalue()).decode()
            base64_images.append(img_data)
        
        return base64_images
    
    def generate_html_report(self, pdf1_path: str, pdf2_path: str,
                           images1_b64: List[str], images2_b64: List[str],
                           differences: Dict) -> str:
        """Generate HTML report with side-by-side comparison."""
        
        total_deletions = len(differences['deletions'])
        total_insertions = len(differences['insertions'])
        total_modifications = len([x for x in differences['modifications'] if x[0] == 'old'])
        total_changes = total_deletions + total_insertions + total_modifications
        
        html_content = f"""
<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>Comparison Report</title>
    <style>
        * {{
            margin: 0;
            padding: 0;
            box-sizing: border-box;
        }}
        
        body {{
            font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, Oxygen, Ubuntu, Cantarell, sans-serif;
            background: #f8fafc;
            color: #1e293b;
            line-height: 1.6;
        }}
        
        .container {{
            max-width: 1400px;
            margin: 0 auto;
            padding: 20px;
        }}
        
        .header {{
            background: white;
            border-radius: 12px;
            padding: 32px;
            margin-bottom: 24px;
            box-shadow: 0 1px 3px rgba(0, 0, 0, 0.1);
            border: 1px solid #e2e8f0;
        }}
        
        .header h1 {{
            font-size: 1.875rem;
            font-weight: 600;
            color: #0f172a;
            margin-bottom: 8px;
        }}
        
        .header p {{
            color: #64748b;
            font-size: 0.875rem;
        }}
        
        .info-section {{
            background: white;
            border-radius: 12px;
            padding: 24px;
            margin-bottom: 24px;
            box-shadow: 0 1px 3px rgba(0, 0, 0, 0.1);
            border: 1px solid #e2e8f0;
        }}
        
        .file-info {{
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 20px;
            margin-bottom: 32px;
        }}
        
        .file-card {{
            background: #f8fafc;
            border: 1px solid #e2e8f0;
            border-radius: 8px;
            padding: 20px;
        }}
        
        .file-card h4 {{
            font-size: 0.875rem;
            font-weight: 600;
            color: #64748b;
            text-transform: uppercase;
            letter-spacing: 0.05em;
            margin-bottom: 12px;
        }}
        
        .file-card p {{
            margin: 6px 0;
            font-size: 0.875rem;
        }}
        
        .file-card strong {{
            color: #374151;
            font-weight: 500;
        }}
        
        .legend {{
            display: flex;
            justify-content: center;
            gap: 24px;
            margin: 32px 0;
            flex-wrap: wrap;
        }}
        
        .legend-item {{
            display: flex;
            align-items: center;
            gap: 8px;
            background: white;
            padding: 12px 16px;
            border-radius: 8px;
            border: 1px solid #e2e8f0;
            font-size: 0.875rem;
            font-weight: 500;
        }}
        
        .legend-color {{
            width: 12px;
            height: 12px;
            border-radius: 3px;
        }}
        
        .summary {{
            background: #f8fafc;
            border: 1px solid #e2e8f0;
            border-radius: 8px;
            padding: 24px;
            text-align: center;
        }}
        
        .summary h3 {{
            font-size: 1.125rem;
            font-weight: 600;
            color: #0f172a;
            margin-bottom: 20px;
        }}
        
        .stats {{
            display: grid;
            grid-template-columns: repeat(auto-fit, minmax(120px, 1fr));
            gap: 16px;
        }}
        
        .stat-item {{
            background: white;
            border: 1px solid #e2e8f0;
            border-radius: 8px;
            padding: 20px 16px;
            text-align: center;
            transition: all 0.2s ease;
        }}
        
        .stat-item:hover {{
            transform: translateY(-2px);
            box-shadow: 0 4px 12px rgba(0, 0, 0, 0.1);
        }}
        
        .stat-number {{
            font-size: 1.875rem;
            font-weight: 700;
            color: #0f172a;
            margin-bottom: 4px;
        }}
        
        .stat-label {{
            color: #64748b;
            font-size: 0.75rem;
            text-transform: uppercase;
            letter-spacing: 0.05em;
            font-weight: 500;
        }}
        
        .comparison-section {{
            display: flex;
            flex-direction: column;
            gap: 24px;
        }}
        
        .page-container {{
            background: white;
            border-radius: 12px;
            overflow: hidden;
            box-shadow: 0 1px 3px rgba(0, 0, 0, 0.1);
            border: 1px solid #e2e8f0;
        }}
        
        .page-header {{
            background: #f8fafc;
            border-bottom: 1px solid #e2e8f0;
            padding: 16px 24px;
            font-weight: 600;
            font-size: 0.875rem;
            color: #374151;
        }}
        
        .page-comparison {{
            display: grid;
            grid-template-columns: 1fr 1fr;
            gap: 1px;
            background: #e2e8f0;
        }}
        
        .page-side {{
            background: white;
            padding: 24px;
            text-align: center;
        }}
        
        .page-side h4 {{
            font-size: 0.75rem;
            font-weight: 600;
            color: #64748b;
            text-transform: uppercase;
            letter-spacing: 0.05em;
            margin-bottom: 16px;
        }}
        
        .page-image {{
            max-width: 100%;
            height: auto;
            border-radius: 8px;
            box-shadow: 0 2px 8px rgba(0, 0, 0, 0.1);
            transition: transform 0.2s ease;
        }}
        
        .page-image:hover {{
            transform: scale(1.02);
        }}
        
        .navigation {{
            position: fixed;
            top: 50%;
            right: 24px;
            transform: translateY(-50%);
            background: white;
            border: 1px solid #e2e8f0;
            border-radius: 12px;
            padding: 16px;
            box-shadow: 0 4px 12px rgba(0, 0, 0, 0.1);
            max-height: 400px;
            overflow-y: auto;
            min-width: 120px;
        }}
        
        .navigation h4 {{
            font-size: 0.75rem;
            font-weight: 600;
            color: #64748b;
            text-transform: uppercase;
            letter-spacing: 0.05em;
            margin-bottom: 12px;
            padding-bottom: 8px;
            border-bottom: 1px solid #e2e8f0;
        }}
        
        .nav-item {{
            display: block;
            padding: 8px 12px;
            text-decoration: none;
            color: #374151;
            border-radius: 6px;
            margin-bottom: 4px;
            transition: all 0.2s ease;
            font-size: 0.875rem;
            font-weight: 500;
        }}
        
        .nav-item:hover {{
            background: #f1f5f9;
            color: #0f172a;
        }}
        
        @media (max-width: 768px) {{
            .container {{
                padding: 16px;
            }}
            
            .file-info {{
                grid-template-columns: 1fr;
                gap: 16px;
            }}
            
            .page-comparison {{
                grid-template-columns: 1fr;
            }}
            
            .navigation {{
                display: none;
            }}
            
            .legend {{
                flex-direction: column;
                align-items: center;
                gap: 12px;
            }}
            
            .stats {{
                grid-template-columns: repeat(2, 1fr);
            }}
            
            .header {{
                padding: 24px;
            }}
            
            .header h1 {{
                font-size: 1.5rem;
            }}
        }}
        
        @media print {{
            .navigation {{
                display: none;
            }}
            
            .page-container {{
                page-break-inside: avoid;
                margin-bottom: 24px;
            }}
        }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>Document Comparison Report</h1>
            <p>Generated on {datetime.now().strftime('%B %d, %Y at %I:%M %p')}</p>
        </div>
        
        <div class="info-section">
            <div class="file-info">
                <div class="file-card">
                    <h4>Original Document</h4>
                    <p>{html.escape(os.path.basename(pdf1_path))}</p>
                </div>
                <div class="file-card">
                    <h4>Modified Document</h4>
                    <p>{html.escape(os.path.basename(pdf2_path))}</p>
                </div>
            </div>
            
            <div class="legend">
                <div class="legend-item">
                    <div class="legend-color" style="background-color: #ef4444;"></div>
                    <span>Deletions</span>
                </div>
                <div class="legend-item">
                    <div class="legend-color" style="background-color: #10b981;"></div>
                    <span>Insertions</span>
                </div>
                <div class="legend-item">
                    <div class="legend-color" style="background-color: #f59e0b;"></div>
                    <span>Modifications</span>
                </div>
            </div>
            
            <div class="summary">
                <h3>Summary</h3>
                <div class="stats">
                    <div class="stat-item">
                        <div class="stat-number">{total_changes}</div>
                        <div class="stat-label">Total Changes</div>
                    </div>
                    <div class="stat-item">
                        <div class="stat-number">{total_deletions}</div>
                        <div class="stat-label">Deletions</div>
                    </div>
                    <div class="stat-item">
                        <div class="stat-number">{total_insertions}</div>
                        <div class="stat-label">Insertions</div>
                    </div>
                    <div class="stat-item">
                        <div class="stat-number">{total_modifications}</div>
                        <div class="stat-label">Modifications</div>
                    </div>
                </div>
            </div>
        </div>
        
        <div class="comparison-section">
"""
        
        # Add page-by-page comparisons
        for page_num in range(len(images1_b64)):
            html_content += f"""
            <div class="page-container" id="page-{page_num + 1}">
                <div class="page-header">
                    Page {page_num + 1}
                </div>
                <div class="page-comparison">
                    <div class="page-side">
                        <h4>Original</h4>
                        <img src="data:image/png;base64,{images1_b64[page_num]}" 
                             alt="Original Page {page_num + 1}" class="page-image">
                    </div>
                    <div class="page-side">
                        <h4>Modified</h4>
                        <img src="data:image/png;base64,{images2_b64[page_num]}" 
                             alt="Modified Page {page_num + 1}" class="page-image">
                    </div>
                </div>
            </div>
            """
        
        # Add navigation
        html_content += """
        </div>
    </div>
    
    <div class="navigation">
        <h4>Pages</h4>
"""
        
        for page_num in range(len(images1_b64)):
            html_content += f"""
        <a href="#page-{page_num + 1}" class="nav-item">Page {page_num + 1}</a>
"""
        
        html_content += """
    </div>
    
    <script>
        // Smooth scrolling for navigation links
        document.querySelectorAll('a[href^="#"]').forEach(anchor => {
            anchor.addEventListener('click', function (e) {
                e.preventDefault();
                const target = document.querySelector(this.getAttribute('href'));
                if (target) {
                    target.scrollIntoView({
                        behavior: 'smooth',
                        block: 'start'
                    });
                }
            });
        });
    </script>
</body>
</html>
"""
        
        return html_content
    
    def compare_pdfs(self, file1_path: str, file2_path: str):
        """Main comparison method for documents (PDF, DOCX, XLSX, PPTX)."""
        print("Starting document comparison...")
        
        # Step 1: Convert documents to PDF if needed
        pdf1_path = self.convert_to_pdf(file1_path)
        pdf2_path = self.convert_to_pdf(file2_path)
        
        # Step 2: Convert PDFs to images
        images1 = self.pdf_to_images(pdf1_path)
        images2 = self.pdf_to_images(pdf2_path)
        
        # Step 3: Extract text blocks with positions
        blocks1 = self.extract_text_blocks(pdf1_path)
        blocks2 = self.extract_text_blocks(pdf2_path)
        
        # Step 4: Pad shorter PDF with empty pages
        images1, images2, blocks1, blocks2 = self.pad_images_and_blocks(
            images1, images2, blocks1, blocks2)
        
        # Step 5: Find text differences
        differences = self.find_text_differences(blocks1, blocks2)
        
        # Step 6: Annotate images with differences
        annotated1, annotated2 = self.annotate_images(images1, images2, differences)
        
        # Step 7: Convert images to base64 for HTML
        images1_b64 = self.save_images_to_base64(annotated1, "original")
        images2_b64 = self.save_images_to_base64(annotated2, "modified")
        
        # Step 8: Generate HTML report
        html_content = self.generate_html_report(
            file1_path, file2_path, images1_b64, images2_b64, differences)
        
        # Step 8: Save HTML report
        report_path = self.comparison_dir / "comparison_report.html"
        with open(report_path, 'w', encoding='utf-8') as f:
            f.write(html_content)
        
        print(f"\n{'='*60}")
        print("DOCUMENT COMPARISON COMPLETE")
        print(f"{'='*60}")
        print(f"Comparison folder: {self.comparison_dir}")
        print(f"Report saved: {report_path}")
        print(f"Total changes found: {len(differences['deletions']) + len(differences['insertions']) + len([x for x in differences['modifications'] if x[0] == 'old'])}")
        print(f"  - Deletions: {len(differences['deletions'])}")
        print(f"  - Insertions: {len(differences['insertions'])}")
        print(f"  - Modifications: {len([x for x in differences['modifications'] if x[0] == 'old'])}")
        print(f"{'='*60}")
        
        return str(report_path)


def main():
    """Main function with CLI interface."""
    parser = argparse.ArgumentParser(
        description="Advanced Document Comparison Tool with visual annotations - supports PDF, DOCX, XLSX, PPTX",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
    python doc_compare.py file1.pdf file2.pdf
    python doc_compare.py document1.docx document2.docx
    python doc_compare.py spreadsheet1.xlsx spreadsheet2.xlsx
    python doc_compare.py presentation1.pptx presentation2.pptx
    python doc_compare.py file1.docx file2.pdf
    
The tool will:
1. Convert documents to PDF format if needed (DOCX, XLSX, PPTX)
2. Convert PDFs to high-resolution images
3. Extract text with position information
4. Pad shorter document with empty pages
5. Find content differences (ignoring layout-only changes)
6. Mark changes on images with colors:
   - Red: Deletions
   - Green: Insertions  
   - Orange: Modifications
7. Generate HTML report in temp/ directory
        """
    )
    
    parser.add_argument('file1', help='Path to the first document file (original) - supports PDF, DOCX, XLSX, PPTX')
    parser.add_argument('file2', help='Path to the second document file (modified) - supports PDF, DOCX, XLSX, PPTX')
    
    args = parser.parse_args()
    
    # Validate input files
    if not Path(args.file1).exists():
        print(f"Error: File '{args.file1}' not found.")
        sys.exit(1)
    
    if not Path(args.file2).exists():
        print(f"Error: File '{args.file2}' not found.")
        sys.exit(1)
    
    try:
        comparator = DocumentComparator()
        report_path = comparator.compare_pdfs(args.file1, args.file2)
        print(f"\nOpen the report in your browser: file://{os.path.abspath(report_path)}")
        
    except Exception as e:
        print(f"Error during comparison: {str(e)}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()